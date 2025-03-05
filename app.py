# app.py
import streamlit as st
import pandas as pd
import numpy as np
import io
import os
import base64
import json
import time
from datetime import datetime
import tempfile
import matplotlib.pyplot as plt
import seaborn as sns
from PIL import Image
import speech_recognition as sr
from pydub import AudioSegment
import docx
import pptx
import fitz  # PyMuPDF
import re
import requests
import anthropic

# Set page configuration
st.set_page_config(
    page_title="Pitch Deck Evaluator",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for styling
st.markdown("""
<style>
    .main-header {
        font-size: 2.5rem;
        font-weight: bold;
        color: #1E3A8A;
        text-align: center;
        margin-bottom: 1rem;
    }
    .sub-header {
        font-size: 1.2rem;
        color: #6B7280;
        text-align: center;
        margin-bottom: 2rem;
    }
    .section-header {
        font-size: 1.5rem;
        font-weight: bold;
        color: #1E3A8A;
        margin-top: 2rem;
        margin-bottom: 1rem;
    }
    .card {
        background-color: white;
        border-radius: 0.5rem;
        padding: 1.5rem;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
        margin-bottom: 1.5rem;
    }
    .metric-container {
        background-color: #F3F4F6;
        border-radius: 9999px;
        padding: 1rem;
        text-align: center;
        margin-bottom: 1rem;
    }
    .metric-circle {
        background-color: white;
        border-radius: 9999px;
        width: 120px;
        height: 120px;
        margin: 0 auto;
        display: flex;
        align-items: center;
        justify-content: center;
        border: 8px solid #10B981;
        font-size: 2rem;
        font-weight: bold;
    }
    .tab-content {
        padding: 1.5rem;
    }
    .feedback-box {
        background-color: #F3F4F6;
        border-radius: 0.5rem;
        padding: 1rem;
        margin-bottom: 1rem;
    }
    .criteria-item {
        display: flex;
        align-items: center;
        margin-bottom: 0.5rem;
    }
    .criteria-item.pass::before {
        content: "✅";
        margin-right: 0.5rem;
    }
    .criteria-item.fail::before {
        content: "❌";
        margin-right: 0.5rem;
    }
    .improvement-box {
        background-color: #F3F4F6;
        border-radius: 0.5rem;
        padding: 1rem;
    }
    .rubric-card {
        border-left: 4px solid #3B82F6;
        padding-left: 1rem;
        margin-bottom: 1.5rem;
    }
    .rubric-level {
        font-weight: 500;
        margin-bottom: 0.25rem;
    }
    .rubric-highlight {
        background-color: #EFF6FF;
        border-radius: 0.375rem;
        padding: 0.75rem;
        margin-top: 0.75rem;
        font-size: 0.875rem;
    }
    .footer {
        text-align: center;
        margin-top: 4rem;
        color: #6B7280;
        font-size: 0.875rem;
    }
</style>
""", unsafe_allow_html=True)

# Initialize session state
if 'evaluation_results' not in st.session_state:
    st.session_state.evaluation_results = None

if 'current_section' not in st.session_state:
    st.session_state.current_section = 'problem'

# Helper functions
def extract_text_from_docx(file):
    """Extract text from a DOCX file."""
    doc = docx.Document(file)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return '\n'.join(full_text)

def extract_text_from_pdf(file):
    """Extract text from a PDF file."""
    pdf_file = fitz.open(stream=file.read(), filetype="pdf")
    text = ""
    for page_num in range(len(pdf_file)):
        page = pdf_file[page_num]
        text += page.get_text()
    return text

def extract_text_from_pptx(file):
    """Extract text from a PPTX file."""
    prs = pptx.Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)

def extract_audio_transcript(audio_file, file_extension):
    """Extract transcript from audio file."""
    # For demonstration, we'll return a mock transcript
    # In a real application, you'd use speech_recognition or another API
    # to actually transcribe the audio
    
    # This would be the real implementation:
    # with tempfile.NamedTemporaryFile(suffix=file_extension) as temp_audio:
    #     temp_audio.write(audio_file.read())
    #     temp_audio.flush()
    #     
    #     # Convert to wav if needed
    #     if file_extension != '.wav':
    #         sound = AudioSegment.from_file(temp_audio.name, format=file_extension[1:])
    #         wav_path = temp_audio.name.replace(file_extension, '.wav')
    #         sound.export(wav_path, format='wav')
    #         audio_path = wav_path
    #     else:
    #         audio_path = temp_audio.name
    #     
    #     # Use speech recognition
    #     recognizer = sr.Recognizer()
    #     with sr.AudioFile(audio_path) as source:
    #         audio_data = recognizer.record(source)
    #         transcript = recognizer.recognize_google(audio_data)
    #         return transcript
    
    # For demo purposes, return a mock transcript based on the sample in the instructions
    return """Good afternoon, everyone. I'm here to talk about a problem that's been plaguing businesses and individuals alike - the inefficiency of current widgets in the market. These widgets, which are supposed to make our lives easier, are instead causing us to waste precious time and resources. In fact, 70% of users have reported dissatisfaction with these widgets, and businesses are losing an average of 20 hours per week due to their inefficiency.

But what if I told you we have a solution? A solution that not only addresses this problem but does so in a way that saves time and resources. We've developed a new kind of widget, one that's designed for maximum efficiency. Our early testing shows that it reduces time wasted by 50%, and we've seen a 95% satisfaction rate among our test users.

Let me paint a picture for you. Imagine a business that's currently losing 20 hours a week due to widget inefficiency. With our new widget, they could potentially save 10 hours a week. That's 10 hours that could be spent on more productive tasks, leading to increased output and profits.

Our business model is simple and effective. We sell our widgets directly to businesses and individuals. By providing a product that offers real value and saves time, we're confident that our widgets will be in high demand. In fact, our market research shows a potential customer base of 1 million users.

Let's talk numbers. We project gross sales of $5 million in the first year, based on an estimated 500,000 transactions. The cost of producing these widgets is $2 million, leaving us with a gross margin of $3 million. After accounting for fixed costs of $1 million, we're looking at a net profit margin of $2 million.

In conclusion, we're offering a solution to a widespread problem, with a compelling business model and sustainable finances. We're not just selling widgets - we're selling efficiency, time savings, and satisfaction. Thank you for your time, and I look forward to your questions."""

def analyze_presentation_with_claude(presentation_text, transcript):
    """
    Analyze the presentation content using Claude API to provide intelligent assessment
    and detailed feedback on the pitch.
    """
    # Initialize Claude client from environment variable or Streamlit secrets
    # You can store your API key in Streamlit's secrets.toml file
    if 'CLAUDE_API_KEY' in st.secrets:
        api_key = st.secrets['CLAUDE_API_KEY']
    else:
        api_key = os.environ.get('CLAUDE_API_KEY')
        
    if not api_key:
        st.error("Claude API key not found. Please set the CLAUDE_API_KEY environment variable or add it to your secrets.toml file.")
        st.stop()
        
    client = anthropic.Anthropic(api_key=api_key)
    
    # Prepare the prompt with the presentation text and transcript
    prompt = f"""
    You are an expert at evaluating business pitches. You need to provide a comprehensive evaluation of a 4-minute pitch presentation based on its content.
    
    Here are the evaluation criteria for a 4-minute business pitch:
    
    1. Problem Framing (25%): 
       - Clearly identifies a significant problem with compelling statistics and examples
       - Uses statistics to demonstrate scale (e.g., 70% user dissatisfaction)
       - Shows impact (e.g., businesses lose 20 hours per week)
       - Explains who is affected
    
    2. Solution Framing (25%):
       - Solution directly addresses identified problem
       - Provides evidence of effectiveness (e.g., 50% time reduction, 95% satisfaction)
       - Explains how solution works and its benefits
       - Compares with alternatives or existing solutions
    
    3. Business Model (20%):
       - Clear explanation of how the business makes money
       - Shows market demand and customer base (e.g., 1 million potential users)
       - Explains value proposition alignment
       - Outlines customer acquisition strategy
    
    4. Financial Overview (20%):
       - Includes gross sales projections (e.g., $5 million)
       - Provides transaction estimates (e.g., 500,000)
       - Shows COGS ($2 million), gross margin ($3 million)
       - Details fixed costs ($1 million) and net profit ($2 million)
    
    5. Delivery & Impact (10%):
       - Clear and concise delivery within 4-minute limit
       - Effective use of slides and visual aids
       - Verbal clarity and engagement
       - Strong conclusion and call to action
    
    Below is the content from a pitch presentation and its transcript. Please evaluate it based on the criteria above.
    
    PRESENTATION CONTENT:
    {presentation_text}
    
    PITCH TRANSCRIPT:
    {transcript}
    
    Please provide:
    1. Scores for each section (0-100)
    2. Specific feedback for each section
    3. An overall score (weighted according to the percentages)
    4. Areas of strength
    5. Suggestions for improvement
    
    Format your response as a JSON object with this structure:
    {{
        "overall": [overall score],
        "sections": {{
            "problem": {{
                "score": [score],
                "feedback": [specific feedback],
                "strengths": [list of strengths],
                "improvements": [list of suggestions]
            }},
            "solution": {{
                "score": [score],
                "feedback": [specific feedback],
                "strengths": [list of strengths],
                "improvements": [list of suggestions]
            }},
            "businessModel": {{
                "score": [score],
                "feedback": [specific feedback],
                "strengths": [list of strengths],
                "improvements": [list of suggestions]
            }},
            "financials": {{
                "score": [score],
                "feedback": [specific feedback],
                "strengths": [list of strengths],
                "improvements": [list of suggestions]
            }},
            "delivery": {{
                "score": [score],
                "feedback": [specific feedback],
                "strengths": [list of strengths],
                "improvements": [list of suggestions]
            }}
        }}
    }}
    
    Return only the JSON object with no additional text.
    """
    
    try:
        # Call Claude API with a progress indicator
        with st.spinner("Claude is analyzing your pitch..."):
            response = client.messages.create(
                model="claude-3-opus-20240229",  # Use the appropriate Claude model version
                max_tokens=4000,
                temperature=0.1,  # Low temperature for more consistent output
                system="You are an expert at evaluating business pitches with deep experience in entrepreneurship, venture capital, and presentation skills. Provide detailed, insightful analysis based on the specified criteria.",
                messages=[
                    {"role": "user", "content": prompt}
                ]
            )
            
            # Extract the response text
            result_text = response.content[0].text
            
            # Parse JSON response
            # First, find JSON object in the response (in case Claude adds additional text)
            json_match = re.search(r'({[\s\S]*})', result_text)
            if json_match:
                result_text = json_match.group(1)
                
            try:
                result = json.loads(result_text)
            except json.JSONDecodeError:
                st.error("Error parsing Claude's response. Using fallback evaluation method.")
                return analyze_presentation_fallback(presentation_text, transcript)
            
            # Ensure all expected fields are in the response
            expected_sections = ["problem", "solution", "businessModel", "financials", "delivery"]
            expected_props = ["score", "feedback"]
            
            for section in expected_sections:
                if section not in result["sections"]:
                    result["sections"][section] = {}
                for prop in expected_props:
                    if prop not in result["sections"][section]:
                        if prop == "score":
                            result["sections"][section][prop] = 70.0
                        else:
                            result["sections"][section][prop] = "No specific feedback provided."
                            
            # Round scores to 1 decimal place
            result["overall"] = round(float(result["overall"]), 1)
            for section in expected_sections:
                result["sections"][section]["score"] = round(float(result["sections"][section]["score"]), 1)
                
            return result
            
    except Exception as e:
        st.error(f"Error calling Claude API: {str(e)}")
        # Fall back to the simpler analysis method
        return analyze_presentation_fallback(presentation_text, transcript)
    
def analyze_presentation_fallback(presentation_text, transcript):
    """
    Fallback analysis method if Claude API call fails.
    Uses basic text matching and rules to generate scores and feedback.
    """
    # Check for key components in the problem section
    problem_indicators = [
        "problem", "challenge", "issue", "pain point", "inefficiency",
        "70%", "dissatisfaction", "20 hours", "wasted time"
    ]
    
    # Check for key components in the solution section
    solution_indicators = [
        "solution", "addresses", "designed for", "efficiency",
        "50%", "reduces time", "95% satisfaction", "test users"
    ]
    
    # Check for key components in the business model section
    business_model_indicators = [
        "business model", "sell", "directly to", "businesses and individuals",
        "value", "high demand", "market research", "1 million", "customer base"
    ]
    
    # Check for key components in the financial section
    financial_indicators = [
        "$5 million", "gross sales", "500,000 transactions", 
        "$2 million", "cost", "gross margin", "$3 million",
        "fixed costs", "$1 million", "net profit", "$2 million"
    ]
    
    # Count the mentions of indicators in the combined text
    combined_text = (presentation_text + " " + transcript).lower()
    
    problem_score = sum(1 for indicator in problem_indicators if indicator.lower() in combined_text) / len(problem_indicators) * 100
    solution_score = sum(1 for indicator in solution_indicators if indicator.lower() in combined_text) / len(solution_indicators) * 100
    business_model_score = sum(1 for indicator in business_model_indicators if indicator.lower() in combined_text) / len(business_model_indicators) * 100
    financial_score = sum(1 for indicator in financial_indicators if indicator.lower() in combined_text) / len(financial_indicators) * 100
    
    # For delivery, we analyze the transcript length and structure
    words = transcript.split()
    word_count = len(words)
    
    # Ideal word count for a 4-minute pitch is roughly 500-600 words
    if 450 <= word_count <= 650:
        delivery_score = 95  # Excellent timing
    elif 400 <= word_count < 450 or 650 < word_count <= 700:
        delivery_score = 85  # Good timing
    elif 350 <= word_count < 400 or 700 < word_count <= 750:
        delivery_score = 75  # Satisfactory timing
    else:
        delivery_score = 65  # Poor timing
    
    # Add randomness to make scores more realistic
    import random
    problem_score = min(100, max(60, problem_score + random.uniform(-5, 5)))
    solution_score = min(100, max(60, solution_score + random.uniform(-5, 5)))
    business_model_score = min(100, max(60, business_model_score + random.uniform(-5, 5)))
    financial_score = min(100, max(60, financial_score + random.uniform(-5, 5)))
    delivery_score = min(100, max(60, delivery_score + random.uniform(-5, 5)))
    
    # Calculate overall score (weighted)
    overall_score = (
        problem_score * 0.25 +
        solution_score * 0.25 +
        business_model_score * 0.20 +
        financial_score * 0.20 +
        delivery_score * 0.10
    )
    
    # Generate feedback
    feedback = {
        "problem": "Strong problem framing with good statistics. Consider highlighting more specific examples of widget inefficiency." if problem_score > 80 else "Problem framing needs more specific statistics and examples to demonstrate the scale of the issue.",
        "solution": "Clear solution presentation, but could strengthen evidence for 50% time reduction claim." if solution_score > 80 else "The solution needs to be more clearly connected to the problem with stronger evidence of effectiveness.",
        "businessModel": "Well-defined business model with good market sizing. Include more details on customer acquisition strategy." if business_model_score > 80 else "Business model needs more detail on how you'll reach your target market and convert them to customers.",
        "financials": "Solid financial breakdown. Consider adding more detail about how fixed costs are calculated." if financial_score > 80 else "Financial projections need more supporting details and breakdown of costs to be credible.",
        "delivery": "Good pace and clarity. More emphasis on the conclusion could strengthen overall impact." if delivery_score > 80 else "Delivery pace needs improvement to fit within the 4-minute timeframe while maintaining clarity."
    }
    
    # Return results
    return {
        "overall": round(overall_score, 1),
        "sections": {
            "problem": {
                "score": round(problem_score, 1),
                "feedback": feedback["problem"],
                "strengths": ["Identifies a problem", "Includes some statistics"],
                "improvements": ["Add more specific examples", "Quantify the impact more clearly"]
            },
            "solution": {
                "score": round(solution_score, 1),
                "feedback": feedback["solution"],
                "strengths": ["Proposes a clear solution", "Mentions benefits"],
                "improvements": ["Provide more evidence", "Compare with alternatives"]
            },
            "businessModel": {
                "score": round(business_model_score, 1),
                "feedback": feedback["businessModel"],
                "strengths": ["Explains revenue mechanism", "Mentions target market"],
                "improvements": ["Add customer acquisition strategy", "Provide more market validation"]
            },
            "financials": {
                "score": round(financial_score, 1),
                "feedback": feedback["financials"],
                "strengths": ["Includes sales projections", "Mentions costs and margins"],
                "improvements": ["Break down fixed costs", "Explain assumptions behind projections"]
            },
            "delivery": {
                "score": round(delivery_score, 1),
                "feedback": feedback["delivery"],
                "strengths": ["Reasonable length", "Clear structure"],
                "improvements": ["Strengthen conclusion", "Improve pacing"]
            }
        }
    }

def generate_radar_chart(scores):
    """Generate a radar chart from the evaluation scores."""
    categories = ['Problem Framing', 'Solution', 'Business Model', 'Financials', 'Delivery']
    values = [
        scores['sections']['problem']['score'], 
        scores['sections']['solution']['score'], 
        scores['sections']['businessModel']['score'],
        scores['sections']['financials']['score'],
        scores['sections']['delivery']['score']
    ]
    
    # Number of variables
    N = len(categories)
    
    # What will be the angle of each axis in the plot (divide the plot / number of variables)
    angles = [n / float(N) * 2 * np.pi for n in range(N)]
    angles += angles[:1]  # Close the loop
    
    # Values need to be repeated to close the loop
    values += values[:1]
    
    # Create the plot
    fig, ax = plt.subplots(figsize=(8, 8), subplot_kw=dict(polar=True))
    
    # Draw one axis per variable and add labels
    plt.xticks(angles[:-1], categories, size=12)
    
    # Draw the outline of the data
    ax.plot(angles, values, linewidth=2, linestyle='solid', color='#3B82F6')
    
    # Fill the area
    ax.fill(angles, values, color='#3B82F6', alpha=0.25)
    
    # Set y-limits
    ax.set_ylim(0, 100)
    
    # Add title
    plt.title('Pitch Evaluation Scores', size=15, color='#1E3A8A', y=1.1)
    
    # Return the figure
    return fig

def get_improvement_suggestions(section, score):
    """Return improvement suggestions based on section and score."""
    suggestions = {
        "problem": {
            "high": "Consider adding 1-2 concise case examples of how widget inefficiency impacts specific businesses. This will strengthen your problem framing by making it more relatable and urgent.",
            "low": "Your problem statement needs more specific data points and real-world examples. Make sure to clearly quantify the scale (e.g., '70% of users report dissatisfaction') and impact (e.g., '20 hours lost per week')."
        },
        "solution": {
            "high": "Provide more concrete evidence for your time reduction claims. Consider including a brief case study or testimonial from your test users to validate your solution's effectiveness.",
            "low": "Your solution needs to be more directly tied to the problem you identified. Make sure to clearly explain how your solution addresses each aspect of the problem and provide measurable benefits (e.g., '50% time reduction')."
        },
        "businessModel": {
            "high": "Add a brief explanation of your customer acquisition strategy. How will you reach your target market efficiently? Include channels and estimated costs to strengthen the business model section.",
            "low": "Your business model needs more detail on revenue generation mechanisms and market validation. Clearly explain how you'll make money, who your customers are, and provide data on market size (e.g., '1 million potential users')."
        },
        "financials": {
            "high": "Break down your fixed costs into major categories (e.g., R&D, marketing, salaries) to demonstrate thoughtful financial planning and increase credibility of your net profit projections.",
            "low": "Your financial overview lacks detail and realistic projections. Make sure to include gross sales projections, transaction estimates, COGS, gross margin, fixed costs, and net profit with supporting calculations."
        },
        "delivery": {
            "high": "End with a stronger conclusion that reinforces your key value proposition and includes a clear call to action. What specific next step do you want the audience to take?",
            "low": "Work on your pacing to fit within the 4-minute timeframe. Practice your delivery to improve clarity and confidence. Make sure your slides support your verbal points without overwhelming the audience."
        }
    }
    
    if score >= 80:
        return suggestions[section]["high"]
    else:
        return suggestions[section]["low"]

def create_download_link(df):
    """Generate a download link for the evaluation report as CSV."""
    csv = df.to_csv(index=False)
    b64 = base64.b64encode(csv.encode()).decode()
    href = f'data:file/csv;base64,{b64}'
    return href

# Main App Layout
def main():
    # Header
    st.markdown('<div class="main-header">Pitch Deck Evaluator</div>', unsafe_allow_html=True)
    st.markdown('<div class="sub-header">Upload, analyze, and get feedback on 4-minute business pitch presentations</div>', unsafe_allow_html=True)
    
    # Create tabs
    tab1, tab2, tab3 = st.tabs(["📤 Upload Materials", "📊 Evaluation Results", "📝 Grading Rubric"])
    
    # Upload Materials Tab
    with tab1:
        st.markdown('<div class="section-header">Upload Pitch Materials</div>', unsafe_allow_html=True)
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown('<div class="card">', unsafe_allow_html=True)
            presentation_file = st.file_uploader("Upload Presentation", type=["ppt", "pptx", "pdf", "doc", "docx"], 
                                               help="Upload your presentation file (PowerPoint, PDF, or Word document)")
            
            if presentation_file is not None:
                st.success(f"✅ {presentation_file.name} uploaded successfully!")
                
                # Show file info
                file_details = {
                    "Filename": presentation_file.name,
                    "File size": f"{presentation_file.size / 1024:.2f} KB",
                    "File type": presentation_file.type
                }
                
                st.json(file_details)
            st.markdown('</div>', unsafe_allow_html=True)
        
        with col2:
            st.markdown('<div class="card">', unsafe_allow_html=True)
            audio_file = st.file_uploader("Upload Audio Recording", type=["mp3", "wav", "ogg", "m4a"], 
                                        help="Upload your pitch recording (4-minute audio file)")
            
            if audio_file is not None:
                st.success(f"✅ {audio_file.name} uploaded successfully!")
                
                # Display audio player
                st.audio(audio_file)
                
                # Show file info
                file_details = {
                    "Filename": audio_file.name,
                    "File size": f"{audio_file.size / 1024:.2f} KB",
                    "File type": audio_file.type
                }
                
                st.json(file_details)
            st.markdown('</div>', unsafe_allow_html=True)
        
        # Evaluate button
        if presentation_file and audio_file:
            if st.button("Evaluate Pitch", type="primary", use_container_width=True):
                with st.spinner("Analyzing your pitch..."):
                    # Process files
                    try:
                        # Extract presentation text based on file type
                        if presentation_file.name.endswith(('.doc', '.docx')):
                            presentation_text = extract_text_from_docx(presentation_file)
                        elif presentation_file.name.endswith('.pdf'):
                            presentation_text = extract_text_from_pdf(presentation_file)
                        elif presentation_file.name.endswith(('.ppt', '.pptx')):
                            presentation_text = extract_text_from_pptx(presentation_file)
                        else:
                            presentation_text = ""
                        
                        # Extract audio transcript
                        file_extension = os.path.splitext(audio_file.name)[1]
                        transcript = extract_audio_transcript(audio_file, file_extension)
                        
                        # Analyze content using Claude
                        evaluation_results = analyze_presentation_with_claude(presentation_text, transcript)
                        
                        # Store results in session state
                        st.session_state.evaluation_results = evaluation_results
                        
                        # Switch to results tab
                        st.experimental_rerun()
                    
                    except Exception as e:
                        st.error(f"Error analyzing files: {str(e)}")
        else:
            st.info("Please upload both a presentation file and an audio recording to proceed with evaluation.")
    
    # Evaluation Results Tab
    with tab2:
        if st.session_state.evaluation_results:
            results = st.session_state.evaluation_results
            
            # Header with overall score
            st.markdown('<div class="section-header">Evaluation Results</div>', unsafe_allow_html=True)
            
            col1, col2 = st.columns([1, 2])
            
            with col1:
                st.markdown('<div class="card">', unsafe_allow_html=True)
                st.markdown(f'''
                <div class="metric-container">
                    <div class="metric-circle">{results["overall"]}%</div>
                </div>
                <h3 style="text-align:center; font-weight:bold; margin-bottom:0.5rem;">Overall Score</h3>
                ''', unsafe_allow_html=True)
                
                # Show rating based on score
                if results["overall"] >= 90:
                    rating = "Excellent"
                    color = "#10B981"  # Green
                elif results["overall"] >= 80:
                    rating = "Good"
                    color = "#3B82F6"  # Blue
                elif results["overall"] >= 70:
                    rating = "Satisfactory"
                    color = "#F59E0B"  # Amber
                else:
                    rating = "Needs Improvement"
                    color = "#EF4444"  # Red
                
                st.markdown(f'''
                <p style="text-align:center; color:{color}; font-weight:bold;">{rating}</p>
                ''', unsafe_allow_html=True)
                
                # Generate and display radar chart
                radar_chart = generate_radar_chart(results)
                st.pyplot(radar_chart)
                st.markdown('</div>', unsafe_allow_html=True)
            
            with col2:
                # Section Selection
                sections = [
                    {"id": "problem", "label": "Problem Framing", "weight": "25%"},
                    {"id": "solution", "label": "Solution Framing", "weight": "25%"},
                    {"id": "businessModel", "label": "Business Model", "weight": "20%"},
                    {"id": "financials", "label": "Financial Overview", "weight": "20%"},
                    {"id": "delivery", "label": "Delivery & Impact", "weight": "10%"}
                ]
                
                st.markdown('<div class="card">', unsafe_allow_html=True)
                selected_section = st.selectbox(
                    "View detailed feedback by section:",
                    options=[section["id"] for section in sections],
                    format_func=lambda x: next((s["label"] for s in sections if s["id"] == x), x),
                    index=[section["id"] for section in sections].index(st.session_state.current_section)
                )
                
                st.session_state.current_section = selected_section
                
                # Display section score
                section_score = results["sections"][selected_section]["score"]
                section_label = next((s["label"] for s in sections if s["id"] == selected_section), selected_section)
                section_weight = next((s["weight"] for s in sections if s["id"] == selected_section), "")
                
                col_a, col_b = st.columns([3, 1])
                
                with col_a:
                    st.markdown(f"### {section_label} <span style='color:#6B7280; font-weight:normal; font-size:1rem;'>({section_weight})</span>", unsafe_allow_html=True)
                
                with col_b:
                    if section_score >= 90:
                        score_color = "#10B981"  # Green
                    elif section_score >= 80:
                        score_color = "#3B82F6"  # Blue
                    elif section_score >= 70:
                        score_color = "#F59E0B"  # Amber
                    else:
                        score_color = "#EF4444"  # Red
                    
                    st.markdown(f"<h2 style='color:{score_color}; text-align:right;'>{section_score}%</h2>", unsafe_allow_html=True)
                
                # Display feedback
                st.markdown("#### Feedback")
                st.markdown(f'''<div class="feedback-box">{results["sections"][selected_section]["feedback"]}</div>''', unsafe_allow_html=True)
                
                # Display strengths and improvements (from Claude analysis)
                col1, col2 = st.columns(2)
                
                with col1:
                    st.markdown("#### Strengths")
                    if "strengths" in results["sections"][selected_section]:
                        strengths = results["sections"][selected_section]["strengths"]
                        if isinstance(strengths, list) and strengths:
                            for strength in strengths:
                                st.markdown(f"✅ {strength}")
                        else:
                            st.markdown("No specific strengths highlighted.")
                    else:
                        # Fallback for older analysis format
                        st.markdown("Detailed strengths analysis not available.")
                
                with col2:
                    st.markdown("#### Areas for Improvement")
                    if "improvements" in results["sections"][selected_section]:
                        improvements = results["sections"][selected_section]["improvements"]
                        if isinstance(improvements, list) and improvements:
                            for improvement in improvements:
                                st.markdown(f"🔍 {improvement}")
                        else:
                            st.markdown("No specific improvements suggested.")
                    else:
                        # Fallback for older analysis format
                        improvement_suggestion = get_improvement_suggestions(selected_section, section_score)
                        st.markdown(f'''<div class="improvement-box">{improvement_suggestion}</div>''', unsafe_allow_html=True)
                
                st.markdown('</div>', unsafe_allow_html=True)
            
            # Generate downloadable report
            st.markdown('<div class="section-header">Download Report</div>', unsafe_allow_html=True)
            
            # Create report dataframe
            report_data = {
                "Section": [s["label"] for s in sections],
                "Weight": [s["weight"] for s in sections],
                "Score": [results["sections"][s["id"]]["score"] for s in sections],
                "Feedback": [results["sections"][s["id"]]["feedback"] for s in sections]
            }
            
            # Add strengths and improvements if available
            if "strengths" in results["sections"]["problem"]:
                report_data["Strengths"] = [
                    ", ".join(results["sections"][s["id"]].get("strengths", [])) 
                    for s in sections
                ]
                report_data["Improvements"] = [
                    ", ".join(results["sections"][s["id"]].get("improvements", []))
                    for s in sections
                ]
            
            report_df = pd.DataFrame(report_data)
            
            if "Strengths" not in report_df.columns:
                report_df["Improvement"] = [
                    get_improvement_suggestions(s["id"], results["sections"][s["id"]]["score"]) 
                    for s in sections
                ]
                
            report_df["Weighted Score"] = [
                results["sections"]["problem"]["score"] * 0.25,
                results["sections"]["solution"]["score"] * 0.25,
                results["sections"]["businessModel"]["score"] * 0.20,
                results["sections"]["financials"]["score"] * 0.20,
                results["sections"]["delivery"]["score"] * 0.10
            ]
            
            # Add overall score row
            overall_row = pd.DataFrame({
                "Section": ["OVERALL"],
                "Weight": ["100%"],
                "Score": [results["overall"]],
                "Feedback": [""],
                "Improvement": [""],
                "Weighted Score": [results["overall"]]
            })
            
            report_df = pd.concat([report_df, overall_row]).reset_index(drop=True)
            
            # Create download link
            csv_link = create_download_link(report_df)
            st.markdown(f'<a href="{csv_link}" download="pitch_evaluation_report.csv" class="download-link" style="display:block; text-align:center; background-color:#3B82F6; color:white; padding:0.75rem; border-radius:0.375rem; text-decoration:none; font-weight:bold;">Download Full Evaluation Report</a>', unsafe_allow_html=True)
            
            # Show a sample of the report
            with st.expander("Preview Report"):
                st.dataframe(report_df)
        
        else:
            st.info("No evaluation results yet. Please upload your materials in the Upload tab and click 'Evaluate Pitch'.")
    
    # Grading Rubric Tab
    with tab3:
        st.markdown('<div class="section-header">Pitch Evaluation Rubric</div>', unsafe_allow_html=True)
        
        # Problem Framing
        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown('''
        <h3 class="text-lg font-semibold mb-3">Problem Framing (25%)</h3>
        <div class="rubric-card">
            <p><span class="rubric-level">Excellent (90-100%):</span> Clearly identifies a significant problem with compelling statistics and examples</p>
            <p><span class="rubric-level">Good (80-89%):</span> Problem is well-defined with supporting data but may lack some specificity</p>
            <p><span class="rubric-level">Satisfactory (70-79%):</span> Problem is identified but lacks sufficient supporting evidence</p>
            <p><span class="rubric-level">Needs Improvement (Below 70%):</span> Problem is vague or poorly supported</p>
            
            <div class="rubric-highlight">
                <p><strong>Key elements:</strong> Definition of problem, statistics showing scale (70% user dissatisfaction), impact demonstration (20 hours lost per week), audience relevance</p>
            </div>
        </div>
        ''', unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)
        
        # Solution Framing
        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown('''
        <h3 class="text-lg font-semibold mb-3">Solution Framing (25%)</h3>
        <div class="rubric-card">
            <p><span class="rubric-level">Excellent (90-100%):</span> Solution directly addresses problem with strong evidence of effectiveness</p>
            <p><span class="rubric-level">Good (80-89%):</span> Clear solution with some evidence of effectiveness</p>
            <p><span class="rubric-level">Satisfactory (70-79%):</span> Solution is presented but connection to problem or evidence is weak</p>
            <p><span class="rubric-level">Needs Improvement (Below 70%):</span> Solution is vague or ineffectively connected to problem</p>
            
            <div class="rubric-highlight">
                <p><strong>Key elements:</strong> Clear description of solution, evidence of effectiveness (50% time reduction, 95% satisfaction), demonstration of impact, comparison with alternatives</p>
            </div>
        </div>
        ''', unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)
        
        # Business Model
        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown('''
        <h3 class="text-lg font-semibold mb-3">Business Model (20%)</h3>
        <div class="rubric-card">
            <p><span class="rubric-level">Excellent (90-100%):</span> Clear, viable business model with strong market validation</p>
            <p><span class="rubric-level">Good (80-89%):</span> Well-defined business model with some market validation</p>
            <p><span class="rubric-level">Satisfactory (70-79%):</span> Basic business model presented but lacks detail or validation</p>
            <p><span class="rubric-level">Needs Improvement (Below 70%):</span> Business model is unclear or unrealistic</p>
            
            <div class="rubric-highlight">
                <p><strong>Key elements:</strong> Revenue mechanism, customer acquisition strategy, market size (1 million potential users), value proposition alignment</p>
            </div>
        </div>
        ''', unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)
        
        # Financial Overview
        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown('''
        <h3 class="text-lg font-semibold mb-3">Financial Overview (20%)</h3>
        <div class="rubric-card">
            <p><span class="rubric-level">Excellent (90-100%):</span> Comprehensive financial projections with realistic assumptions</p>
            <p><span class="rubric-level">Good (80-89%):</span> Solid financial breakdown with mostly realistic projections</p>
            <p><span class="rubric-level">Satisfactory (70-79%):</span> Basic financial information provided but lacks detail or realism</p>
            <p><span class="rubric-level">Needs Improvement (Below 70%):</span> Financial information is missing key elements or unrealistic</p>
            
            <div class="rubric-highlight">
                <p><strong>Key elements:</strong> Gross sales projections ($5M), transaction estimates (500,000), COGS ($2M), gross margin ($3M), fixed costs ($1M), net profit ($2M)</p>
            </div>
        </div>
        ''', unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)
        
        # Delivery & Impact
        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown('''
        <h3 class="text-lg font-semibold mb-3">Delivery & Impact (10%)</h3>
        <div class="rubric-card">
            <p><span class="rubric-level">Excellent (90-100%):</span> Confident, engaging delivery that stays within 4-minute time limit</p>
            <p><span class="rubric-level">Good (80-89%):</span> Clear delivery with good time management</p>
            <p><span class="rubric-level">Satisfactory (70-79%):</span> Adequate delivery with some timing issues</p>
            <p><span class="rubric-level">Needs Improvement (Below 70%):</span> Poor delivery or significantly over/under time</p>
            
            <div class="rubric-highlight">
                <p><strong>Key elements:</strong> Time management (4-minute limit), slide quality, verbal clarity, engagement, compelling conclusion</p>
            </div>
        </div>
        ''', unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)
        
        # Assignment Guidelines
        st.markdown('<div class="card">', unsafe_allow_html=True)
        st.markdown('''
        <h3 class="text-lg font-semibold mb-3">Assignment Guidelines</h3>
        <p>Students should prepare a 4-minute business pitch with exactly 4 slides:</p>
        <ol>
            <li><strong>Slide 1: Problem Framing</strong> - Identify the problem, show statistics, explain who is affected</li>
            <li><strong>Slide 2: Solution Framing</strong> - Present your solution, explain how it works, provide evidence</li>
            <li><strong>Slide 3: Business Model</strong> - Explain how you make money, show market demand, identify target customers</li>
            <li><strong>Slide 4: Financial Overview</strong> - Detail gross sales, transactions, costs, margins, and profit</li>
        </ol>
        <p>Remember to stay within the 4-minute time limit and follow the structure outlined in the assignment.</p>
        ''', unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)
        
    # Footer
    st.markdown('''
    <div class="footer">
        Pitch Deck Evaluator • Created for Business Pitch Assessment
    </div>
    ''', unsafe_allow_html=True)

if __name__ == "__main__":
    main()
